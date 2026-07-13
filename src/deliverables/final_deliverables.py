"""Generate final project documentation and presentation artifacts."""

from __future__ import annotations

import json
import sqlite3
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path

import pandas as pd
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


PROJECT_ROOT = Path(__file__).resolve().parents[2]
REPORTS_DIR = PROJECT_ROOT / "reports"
FIGURES_DIR = REPORTS_DIR / "figures"
DOCS_DIR = PROJECT_ROOT / "docs" / "final_deliverables"
DATABASE_PATH = PROJECT_ROOT / "data" / "database" / "manufacturing_copilot.db"


@dataclass(frozen=True)
class ProjectMetrics:
    train_products: int
    train_failures: int
    failure_rate_pct: float
    test_products: int
    test_alerts: int
    validation_rows: int
    validation_failures: int
    best_model: str
    mcc: float
    precision: float
    recall: float
    f1: float
    pr_auc: float
    product_families: int
    highest_family: int
    highest_family_failure_rate_pct: float
    top_station: str
    top_station_failure_rate_pct: float
    top_bottleneck: str
    top_bottleneck_score: float
    top_critical_node: str
    kg_nodes: int
    kg_edges: int
    copilot_rows: int


def read_csv(path: Path) -> pd.DataFrame:
    if not path.exists():
        raise FileNotFoundError(path)
    return pd.read_csv(path)


def collect_metrics() -> ProjectMetrics:
    model_metrics = read_csv(REPORTS_DIR / "phase6_model_comparison_metrics.csv")
    best = model_metrics.sort_values("rank").iloc[0]
    families = read_csv(REPORTS_DIR / "phase5_product_family_failure_rates.csv")
    top_family = families.sort_values("failure_rate_pct", ascending=False).iloc[0]
    stations = read_csv(REPORTS_DIR / "phase3_station_failure_rates.csv")
    top_station = stations.sort_values("failure_rate_pct", ascending=False).iloc[0]
    bottlenecks = read_csv(REPORTS_DIR / "phase8_bottleneck_scores.csv")
    top_bottleneck = bottlenecks.sort_values("bottleneck_rank").iloc[0]
    critical_nodes = read_csv(REPORTS_DIR / "phase9_critical_nodes.csv")
    top_node = critical_nodes.sort_values("critical_rank").iloc[0]
    kg_nodes = len(read_csv(REPORTS_DIR / "phase9_knowledge_graph_nodes.csv"))
    kg_edges = len(read_csv(REPORTS_DIR / "phase9_knowledge_graph_edges.csv"))

    with sqlite3.connect(DATABASE_PATH) as connection:
        train_summary = pd.read_sql_query(
            """
            SELECT labeled_count AS train_products, failure_count AS train_failures,
                   failure_rate_pct
            FROM throughput_efficiency WHERE split='train'
            """,
            connection,
        ).iloc[0]
        test_summary = pd.read_sql_query(
            """
            SELECT COUNT(1) AS test_products,
                   SUM(predicted_failure) AS test_alerts
            FROM product_predictions WHERE split='test'
            """,
            connection,
        ).iloc[0]
        validation_summary = pd.read_sql_query(
            """
            SELECT COUNT(1) AS validation_rows,
                   SUM(actual_response) AS validation_failures
            FROM product_predictions WHERE split='validation'
            """,
            connection,
        ).iloc[0]
        copilot_rows = pd.read_sql_query(
            "SELECT COUNT(1) AS rows FROM product_predictions",
            connection,
        ).iloc[0]["rows"]

    return ProjectMetrics(
        train_products=int(train_summary["train_products"]),
        train_failures=int(train_summary["train_failures"]),
        failure_rate_pct=float(train_summary["failure_rate_pct"]),
        test_products=int(test_summary["test_products"]),
        test_alerts=int(test_summary["test_alerts"]),
        validation_rows=int(validation_summary["validation_rows"]),
        validation_failures=int(validation_summary["validation_failures"]),
        best_model=str(best["model"]),
        mcc=float(best["mcc"]),
        precision=float(best["precision"]),
        recall=float(best["recall"]),
        f1=float(best["f1"]),
        pr_auc=float(best["pr_auc"]),
        product_families=int(families["final_product_family"].nunique()),
        highest_family=int(top_family["final_product_family"]),
        highest_family_failure_rate_pct=float(top_family["failure_rate_pct"]),
        top_station=str(top_station["station_key"]),
        top_station_failure_rate_pct=float(top_station["failure_rate_pct"]),
        top_bottleneck=str(top_bottleneck["station"]),
        top_bottleneck_score=float(top_bottleneck["bottleneck_score"]),
        top_critical_node=str(top_node["station"]),
        kg_nodes=kg_nodes,
        kg_edges=kg_edges,
        copilot_rows=int(copilot_rows),
    )


def deliverable_map() -> list[dict[str, object]]:
    rows = [
        {
            "id": 68,
            "deliverable": "Production Failure Prediction Model",
            "primary_artifacts": [
                "models/phase6_best_model.joblib",
                "src/data/phase6_predictive_failure_modeling.py",
                "notebooks/phase6_predictive_failure_modeling.ipynb",
                "reports/phase6_predictive_failure_modeling_report.md",
                "data/processed/phase6_test_predictions.csv",
            ],
            "status": "Complete",
        },
        {
            "id": 69,
            "deliverable": "Product Family Segmentation Module",
            "primary_artifacts": [
                "src/data/phase5_product_family_discovery.py",
                "notebooks/phase5_product_family_discovery.ipynb",
                "reports/phase5_product_family_discovery_report.md",
                "data/processed/phase5_train_product_families.csv",
                "data/processed/phase5_test_product_families.csv",
            ],
            "status": "Complete",
        },
        {
            "id": 70,
            "deliverable": "Process Mining Engine",
            "primary_artifacts": [
                "src/data/phase8_process_mining_bottleneck_analysis.py",
                "notebooks/phase8_process_mining_bottleneck_analysis.ipynb",
                "reports/phase8_process_mining_bottleneck_report.md",
                "reports/phase8_process_map_edges.csv",
                "reports/phase8_bottleneck_scores.csv",
            ],
            "status": "Complete",
        },
        {
            "id": 71,
            "deliverable": "Root Cause Analysis Engine",
            "primary_artifacts": [
                "src/data/phase7_root_cause_analysis.py",
                "notebooks/phase7_root_cause_analysis.ipynb",
                "reports/phase7_root_cause_analysis_report.md",
                "reports/phase7_top_failure_drivers.csv",
                "reports/phase7_station_root_cause_report.csv",
            ],
            "status": "Complete",
        },
        {
            "id": 72,
            "deliverable": "Knowledge Graph",
            "primary_artifacts": [
                "src/data/phase9_knowledge_graph.py",
                "notebooks/phase9_knowledge_graph.ipynb",
                "reports/phase9_manufacturing_knowledge_graph.graphml",
                "reports/phase9_knowledge_graph_report.html",
                "reports/phase9_critical_nodes.csv",
            ],
            "status": "Complete",
        },
        {
            "id": 73,
            "deliverable": "Manufacturing Copilot",
            "primary_artifacts": [
                "app/phase11_manufacturing_copilot.py",
                "src/copilot/query_engine.py",
                "src/data/phase11_manufacturing_copilot.py",
                "data/database/manufacturing_copilot.db",
                "notebooks/phase11_manufacturing_copilot.ipynb",
            ],
            "status": "Complete",
        },
        {
            "id": 74,
            "deliverable": "Executive Dashboard",
            "primary_artifacts": [
                "app/phase12_executive_dashboard.py",
                "src/dashboard/business_impact.py",
                "src/data/phase12_executive_dashboard.py",
                "notebooks/phase12_executive_dashboard.ipynb",
                "reports/phase12_executive_dashboard_report.md",
            ],
            "status": "Complete",
        },
        {
            "id": 75,
            "deliverable": "Project Documentation and Presentation",
            "primary_artifacts": [
                "docs/final_deliverables/Bosch_Production_Line_Performance_Research_Report.docx",
                "docs/final_deliverables/Bosch_Production_Line_Performance_Research_Report.md",
                "docs/final_deliverables/Bosch_Production_Line_Performance_Presentation.pptx",
                "docs/final_deliverables/final_deliverables_index.md",
                "docs/final_deliverables/final_deliverables_manifest.json",
            ],
            "status": "Complete",
        },
    ]
    return rows


def report_markdown(metrics: ProjectMetrics) -> str:
    return f"""# Manufacturing Failure Prediction and Process Intelligence for Bosch Production-Line Data

**Author:** Nexturn Manufacturing Analytics Use Case  
**Document type:** IEEE-style technical research report  
**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M')}  

## Abstract

This project develops an end-to-end manufacturing analytics system using the Bosch Production Line Performance dataset. The system combines failure prediction, product-family segmentation, process mining, root-cause analysis, knowledge-graph analytics, anomaly detection, a manufacturing copilot, and an executive dashboard. The production-safe model is a LightGBM classifier trained from raw numeric, categorical, and date-derived features. On the held-out validation set, the selected model achieved MCC={metrics.mcc:.3f}, precision={metrics.precision:.3f}, recall={metrics.recall:.3f}, F1={metrics.f1:.3f}, and PR-AUC={metrics.pr_auc:.3f}. The analysis identifies {metrics.top_station} as the highest observed risk station with {metrics.top_station_failure_rate_pct:.3f}% failure rate, {metrics.top_bottleneck} as the top bottleneck, and {metrics.top_critical_node} as the top knowledge-graph critical node. The final system is packaged as reusable Python scripts, notebooks, model artifacts, a SQLite-backed copilot, and a Streamlit executive dashboard.

**Keywords:** manufacturing analytics, failure prediction, process mining, root-cause analysis, knowledge graph, SHAP, LightGBM, Bosch production-line data.

## I. Introduction

Modern production lines generate high-dimensional sensor, timestamp, and routing data. The business objective is not only to classify likely failures, but also to explain where quality risk concentrates, how products flow through the process, and which operational interventions should be prioritized. This project converts the Bosch dataset into a reusable manufacturing analytics architecture for future client engagements.

The implemented system covers the complete project roadmap: project setup, data understanding, exploratory analysis, feature engineering, product-family discovery, predictive modeling, root-cause analysis, process mining, knowledge graph construction, advanced AI experiments, a manufacturing copilot, and an executive dashboard.

## II. Dataset and Scope

The raw Bosch data includes numeric, categorical, and date files for train and test populations. The training numeric file contains the target variable `Response`; test labels are unavailable. The project uses raw CSV data for modeling and analysis, while selected processed files are created as reproducible phase outputs.

- Labeled train products: {metrics.train_products:,}
- Train failures: {metrics.train_failures:,}
- Historical failure rate: {metrics.failure_rate_pct:.3f}%
- Test products scored: {metrics.test_products:,}
- Validation products used for model assessment: {metrics.validation_rows:,}
- Validation failures: {metrics.validation_failures:,}

The Kaggle leaderboard-style leakage features are documented separately as research-only. They are not used for production-safe modeling, root-cause explanation, dashboards, or copilot outputs because nearby known labels would not be available in live deployment.

## III. Methodology

### A. Data Quality and Metadata Engineering

Phase 1 loaded raw numeric, categorical, date, and target datasets and documented dataset sizes, feature counts, missing values, and file-level quality checks. Phase 2 parsed feature names into production line, station, and feature identifiers, producing station metadata and completeness summaries.

### B. Exploratory Analysis

Phase 3 calculated failure rates by line and station, analyzed time-derived failure patterns, generated correlation and distribution reports, and used one-hot encoded categorical subsets where required for EDA. The most important station-level observation was that {metrics.top_station} showed a materially elevated failure rate relative to the overall training baseline.

### C. Feature Engineering

Phase 4 constructed product-level timing and path features including start time, end time, cycle time, processing duration, waiting time, station counts, line counts, path complexity, and line-level aggregates. These features created a common interface for predictive modeling, process mining, and dashboarding.

### D. Product Family Discovery

Phase 5 created station-presence matrices and clustered products using KMeans, DBSCAN, and hierarchical clustering. The final segmentation produced {metrics.product_families} product families. Family {metrics.highest_family} had the highest observed failure rate at {metrics.highest_family_failure_rate_pct:.3f}%, indicating that product routes and family mix are meaningful risk segmentation dimensions.

### E. Predictive Modeling

Phase 6 compared logistic regression, random forest, XGBoost, LightGBM, and CatBoost. The selected production-safe model is {metrics.best_model}. It achieved MCC={metrics.mcc:.3f}, precision={metrics.precision:.3f}, recall={metrics.recall:.3f}, F1={metrics.f1:.3f}, and PR-AUC={metrics.pr_auc:.3f}. Later improvement experiments did not outperform the Phase 6 baseline, so the original LightGBM model remains the official failure prediction model.

### F. Root-Cause Analysis

Phase 7 generated model feature importance and SHAP explanations. The strongest global drivers included timing features, line timing features, station-specific numeric measurements, and missingness indicators. SHAP outputs are treated as model explanations and engineering investigation priorities, not as proof of physical causality.

### G. Process Mining and Bottleneck Analysis

Phase 8 reconstructed production process maps from raw date features, calculated waiting-time metrics, identified bottlenecks, and measured throughput efficiency proxies. {metrics.top_bottleneck} ranked first with bottleneck score {metrics.top_bottleneck_score:.2f}. The throughput efficiency metric is a relative timestamp-derived proxy and must not be presented as formal OEE.

### H. Knowledge Graph

Phase 9 created a Line-to-Station-to-Feature-to-Failure graph, built station relationship edges, calculated centrality metrics, identified candidate failure propagation routes, and highlighted critical nodes. The graph contains {metrics.kg_nodes:,} nodes and {metrics.kg_edges:,} edges. {metrics.top_critical_node} is the top critical node by the combined graph and operational priority score.

### I. Advanced AI Experiments

Phase 10 evaluated Isolation Forest, MLP reconstruction anomaly scoring, graph message-passing station risk, and failure trajectory prediction. These methods were useful as diagnostic add-ons but did not replace the production-safe LightGBM model.

### J. Manufacturing Copilot and Executive Dashboard

Phase 11 stores model outputs and analytics artifacts in SQLite, then exposes them through a Streamlit manufacturing copilot. The copilot database contains {metrics.copilot_rows:,} product prediction rows. Phase 12 adds an executive dashboard with KPI cards, relative failure trends, station heatmaps, bottleneck analytics, SHAP summaries, and scenario-based business impact calculations.

## IV. Key Results

### A. The official production-safe model is accurate enough for prioritization, not autonomous disposition.

The LightGBM model achieved MCC={metrics.mcc:.3f} and precision={metrics.precision:.3f}. This supports prioritizing inspection and engineering review. It should not be used as an autonomous scrap/release decision without plant-specific validation, monitoring, and operating procedures.

### B. Quality risk is concentrated by station and path.

{metrics.top_station} is the highest observed risk station at {metrics.top_station_failure_rate_pct:.3f}% failure rate. Product family {metrics.highest_family} has elevated risk at {metrics.highest_family_failure_rate_pct:.3f}%, showing that route and family segmentation add operational context beyond a single global model score.

### C. Bottlenecks and critical nodes are not identical.

{metrics.top_bottleneck} is the top bottleneck, while {metrics.top_critical_node} is the top knowledge-graph critical node. This distinction matters: bottleneck score emphasizes waiting and flow constraint, while graph criticality combines network position, failure rate, bottleneck evidence, and root-cause importance.

### D. Executive impact must be scenario-based.

The project includes a business-impact calculator using production volume, historical failure rate, model alert rate, validation precision, cost per failure, cost per alert review, and intervention effectiveness. The calculator is intentionally assumption-driven because the Kaggle dataset does not include plant-specific financial outcomes or measured intervention effects.

## V. System Architecture

The final architecture has five layers:

1. **Raw data and metadata layer:** Bosch raw CSV files, feature parsing, station metadata, line metadata, completeness metrics.
2. **Analytical feature layer:** timing, waiting, station count, path complexity, product-family labels, graph and trajectory features.
3. **Model layer:** production-safe LightGBM model plus benchmark models and research-only leaderboard variant.
4. **Diagnostic intelligence layer:** SHAP explanations, bottleneck scores, knowledge graph, anomaly scores, failure trajectory scores.
5. **Decision interface layer:** SQLite database, Manufacturing Copilot Streamlit app, Executive Dashboard Streamlit app, reports, notebooks, and presentation.

## VI. Deliverable Inventory

| ID | Deliverable | Status | Primary artifact |
|---:|---|---|---|
| 68 | Production Failure Prediction Model | Complete | `models/phase6_best_model.joblib` |
| 69 | Product Family Segmentation Module | Complete | `src/data/phase5_product_family_discovery.py` |
| 70 | Process Mining Engine | Complete | `src/data/phase8_process_mining_bottleneck_analysis.py` |
| 71 | Root Cause Analysis Engine | Complete | `src/data/phase7_root_cause_analysis.py` |
| 72 | Knowledge Graph | Complete | `reports/phase9_manufacturing_knowledge_graph.graphml` |
| 73 | Manufacturing Copilot | Complete | `app/phase11_manufacturing_copilot.py` |
| 74 | Executive Dashboard | Complete | `app/phase12_executive_dashboard.py` |
| 75 | Documentation and Presentation | Complete | `docs/final_deliverables/` |

## VII. Limitations and Robustness Notes

- Test labels are unavailable, so test predictions cannot be evaluated as measured failures.
- Kaggle competition leaderboard leakage strategies are excluded from production deliverables.
- SHAP and graph routes provide diagnostic associations, not causal proof.
- Bosch relative timestamps do not correspond to real calendar dates.
- Financial impact is scenario-based and requires client-specific cost validation.
- A production deployment would require retraining governance, feature drift monitoring, alert workflow design, security review, and integration with plant systems.

## VIII. Recommended Next Steps

1. Convert the current local pipeline into scheduled production jobs for a future client dataset.
2. Replace competition files with governed plant data sources and stable data contracts.
3. Validate the model with forward-time splits and post-alert quality outcomes.
4. Connect the copilot and dashboard to role-based access control.
5. Add monitoring for feature drift, prediction drift, alert precision, recall, and intervention outcomes.
6. Work with manufacturing engineers to confirm root causes using maintenance, sensor, routing, and inspection records.

## IX. Conclusion

The project delivers a complete manufacturing analytics use case rather than a single classifier. The selected production-safe model provides risk prioritization, while segmentation, process mining, root-cause analysis, graph analytics, copilot querying, and executive reporting make the results interpretable and operationally usable. The system is suitable as a reusable reference architecture for future production-line quality analytics projects, subject to client-specific validation and deployment controls.

## References

[1] Bosch Production Line Performance, Kaggle competition dataset.  
[2] LightGBM: A Highly Efficient Gradient Boosting Decision Tree.  
[3] SHAP: A Unified Approach to Interpreting Model Predictions.  
[4] Process mining and event-log analysis methods for production systems.  
[5] Network centrality and graph analytics methods for manufacturing process graphs.
"""


def write_markdown(metrics: ProjectMetrics) -> Path:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    path = DOCS_DIR / "Bosch_Production_Line_Performance_Research_Report.md"
    path.write_text(report_markdown(metrics), encoding="utf-8")
    return path


def set_normal_style(document: Document) -> None:
    styles = document.styles
    normal = styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = Pt(10)
    for section in document.sections:
        section.top_margin = Inches(0.75)
        section.bottom_margin = Inches(0.75)
        section.left_margin = Inches(0.7)
        section.right_margin = Inches(0.7)


def add_docx_heading(document: Document, text: str, level: int = 1) -> None:
    paragraph = document.add_heading(text, level=level)
    paragraph.style.font.name = "Times New Roman"


def add_docx_image(document: Document, path: Path, caption: str) -> None:
    if not path.exists():
        return
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run()
    run.add_picture(str(path), width=Inches(5.8))
    caption_paragraph = document.add_paragraph(caption)
    caption_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    caption_paragraph.runs[0].italic = True


def write_docx(metrics: ProjectMetrics) -> Path:
    document = Document()
    set_normal_style(document)
    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title.add_run("Manufacturing Failure Prediction and Process Intelligence for Bosch Production-Line Data")
    title_run.bold = True
    title_run.font.size = Pt(16)
    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.add_run("IEEE-style technical research report").italic = True

    add_docx_heading(document, "Abstract", 1)
    document.add_paragraph(
        "This report documents an end-to-end manufacturing analytics system for the Bosch "
        "Production Line Performance dataset. The system includes failure prediction, "
        "product-family segmentation, process mining, root-cause analysis, knowledge-graph "
        "analytics, advanced AI diagnostics, a manufacturing copilot, and an executive dashboard."
    )
    document.add_paragraph(
        f"The selected production-safe model is {metrics.best_model}, with validation "
        f"MCC={metrics.mcc:.3f}, precision={metrics.precision:.3f}, recall={metrics.recall:.3f}, "
        f"F1={metrics.f1:.3f}, and PR-AUC={metrics.pr_auc:.3f}."
    )
    document.add_paragraph(
        "Keywords: manufacturing analytics, failure prediction, process mining, root-cause analysis, knowledge graph, SHAP, LightGBM."
    )

    add_docx_heading(document, "I. Introduction", 1)
    document.add_paragraph(
        "The objective is to move beyond a competition classifier and create a reusable "
        "manufacturing intelligence system that predicts failures, explains likely drivers, "
        "maps product flow, and supports executive decision-making."
    )

    add_docx_heading(document, "II. Dataset and Scope", 1)
    for item in [
        f"Labeled train products: {metrics.train_products:,}",
        f"Train failures: {metrics.train_failures:,}",
        f"Historical failure rate: {metrics.failure_rate_pct:.3f}%",
        f"Test products scored: {metrics.test_products:,}",
        f"Validation rows: {metrics.validation_rows:,}",
    ]:
        document.add_paragraph(item, style="List Bullet")

    add_docx_heading(document, "III. Methodology", 1)
    methodology = [
        "Feature names were parsed into production line, station, and feature identifiers.",
        "Timing, waiting, cycle-time, station-count, path-complexity, and line-level aggregate features were engineered.",
        "Product families were discovered from station-presence patterns using clustering methods.",
        "Failure models were benchmarked across logistic regression, random forest, XGBoost, LightGBM, and CatBoost.",
        "Root-cause diagnostics were produced from feature importance and SHAP explanations.",
        "Process mining reconstructed station transitions, bottlenecks, critical paths, and throughput proxies.",
        "A knowledge graph linked lines, stations, features, and failure evidence.",
        "A SQLite-backed copilot and Streamlit dashboards operationalized the outputs.",
    ]
    for item in methodology:
        document.add_paragraph(item, style="List Number")

    add_docx_heading(document, "IV. Results", 1)
    document.add_paragraph(
        f"The official production-safe model remains {metrics.best_model}. It achieved "
        f"MCC={metrics.mcc:.3f}, precision={metrics.precision:.3f}, recall={metrics.recall:.3f}, "
        f"F1={metrics.f1:.3f}, and PR-AUC={metrics.pr_auc:.3f}."
    )
    document.add_paragraph(
        f"The highest-risk station is {metrics.top_station} with {metrics.top_station_failure_rate_pct:.3f}% "
        f"failure rate. The highest-risk product family is family {metrics.highest_family} with "
        f"{metrics.highest_family_failure_rate_pct:.3f}% failure rate."
    )
    document.add_paragraph(
        f"{metrics.top_bottleneck} is the top bottleneck, while {metrics.top_critical_node} is the "
        "highest-ranked critical knowledge-graph node."
    )

    add_docx_image(
        document,
        FIGURES_DIR / "phase6_improvement_mcc_comparison.png",
        "Fig. 1. Production-safe model comparison and improvement experiments.",
    )
    add_docx_image(
        document,
        FIGURES_DIR / "phase7_top_shap_drivers.png",
        "Fig. 2. Global SHAP failure drivers.",
    )
    add_docx_image(
        document,
        FIGURES_DIR / "phase8_bottleneck_scores.png",
        "Fig. 3. Top process bottlenecks.",
    )
    add_docx_image(
        document,
        FIGURES_DIR / "phase9_critical_nodes.png",
        "Fig. 4. Critical manufacturing knowledge-graph nodes.",
    )

    add_docx_heading(document, "V. Deliverables", 1)
    table = document.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    for idx, text in enumerate(["ID", "Deliverable", "Status", "Primary Artifact"]):
        table.rows[0].cells[idx].text = text
    for row in deliverable_map():
        cells = table.add_row().cells
        cells[0].text = str(row["id"])
        cells[1].text = str(row["deliverable"])
        cells[2].text = str(row["status"])
        cells[3].text = str(row["primary_artifacts"][0])

    add_docx_heading(document, "VI. Limitations", 1)
    for item in [
        "Test labels are unavailable, so test predictions cannot be measured as observed failures.",
        "Leaderboard leakage strategies are excluded from production deliverables.",
        "SHAP and graph routes are diagnostic associations, not causal proof.",
        "Financial impact is scenario-based and requires client-specific validation.",
        "Relative Bosch timestamps do not directly map to calendar dates.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    add_docx_heading(document, "VII. Conclusion", 1)
    document.add_paragraph(
        "The project delivers a complete manufacturing analytics reference architecture for "
        "future production-line quality use cases. It combines model accuracy, interpretability, "
        "process intelligence, graph analytics, and decision interfaces while preserving clear "
        "production-safety boundaries."
    )

    add_docx_heading(document, "References", 1)
    for ref in [
        "[1] Bosch Production Line Performance, Kaggle competition dataset.",
        "[2] LightGBM: A Highly Efficient Gradient Boosting Decision Tree.",
        "[3] SHAP: A Unified Approach to Interpreting Model Predictions.",
        "[4] Process mining and event-log analysis methods for production systems.",
        "[5] Network centrality and graph analytics methods for manufacturing process graphs.",
    ]:
        document.add_paragraph(ref)

    path = DOCS_DIR / "Bosch_Production_Line_Performance_Research_Report.docx"
    document.save(path)
    return path


def add_slide_title(slide, title: str, subtitle: str | None = None) -> None:
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = PptPt(30)
    if subtitle:
        box = slide.shapes.add_textbox(PptInches(0.7), PptInches(1.45), PptInches(12), PptInches(0.45))
        box.text_frame.text = subtitle
        box.text_frame.paragraphs[0].font.size = PptPt(14)


def add_bullets(slide, bullets: list[str], left=0.8, top=1.55, width=11.8, height=4.8) -> None:
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    tf = box.text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = PptPt(18)


def add_picture(slide, image: Path, left=6.7, top=1.55, width=5.7) -> None:
    if image.exists():
        slide.shapes.add_picture(str(image), PptInches(left), PptInches(top), width=PptInches(width))


def write_presentation(metrics: ProjectMetrics) -> Path:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[5]

    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Bosch Production Line Performance"
    slide.placeholders[1].text = "Final Deliverables: Predictive Quality, Process Intelligence, Copilot, and Executive Dashboard"

    slides = [
        (
            "Project Objective",
            [
                "Build a reusable manufacturing analytics reference architecture.",
                "Predict product failure risk from raw numeric, categorical, and date data.",
                "Explain risk through SHAP, station evidence, process bottlenecks, and graph analytics.",
                "Deliver usable interfaces: Manufacturing Copilot and Executive Dashboard.",
            ],
            None,
        ),
        (
            "Data and Scope",
            [
                f"Labeled train products: {metrics.train_products:,}; failures: {metrics.train_failures:,}.",
                f"Historical failure rate: {metrics.failure_rate_pct:.3f}%.",
                f"Test products scored: {metrics.test_products:,}; predicted alerts: {metrics.test_alerts:,}.",
                "Raw Bosch CSVs remain the source. Processed outputs are reproducible phase artifacts.",
            ],
            None,
        ),
        (
            "Production Failure Prediction Model",
            [
                f"Selected model: {metrics.best_model}.",
                f"Validation MCC {metrics.mcc:.3f}; precision {metrics.precision:.3f}; recall {metrics.recall:.3f}.",
                "Improvement experiments did not beat the Phase 6 LightGBM baseline.",
                "Leaderboard leakage model remains research-only and excluded from production deliverables.",
            ],
            FIGURES_DIR / "phase6_improvement_mcc_comparison.png",
        ),
        (
            "Product Family Segmentation",
            [
                f"Final segmentation contains {metrics.product_families} product families.",
                f"Highest-risk family: {metrics.highest_family} at {metrics.highest_family_failure_rate_pct:.3f}% failure rate.",
                "Families are based on station-presence and flow-path structure.",
                "Segmentation supports targeted monitoring by route and product mix.",
            ],
            None,
        ),
        (
            "Process Mining Engine",
            [
                f"Top bottleneck: {metrics.top_bottleneck} with score {metrics.top_bottleneck_score:.2f}.",
                "Outputs include process-map edges, station waiting times, bottleneck scores, and critical paths.",
                "Throughput efficiency is a relative timestamp proxy, not formal OEE.",
            ],
            FIGURES_DIR / "phase8_bottleneck_scores.png",
        ),
        (
            "Root Cause Analysis Engine",
            [
                "SHAP identifies timing, line-level, missingness, and station-measurement drivers.",
                f"Highest observed station risk: {metrics.top_station} at {metrics.top_station_failure_rate_pct:.3f}%.",
                "Recommendations are engineering investigation priorities, not causal proof.",
            ],
            FIGURES_DIR / "phase7_top_shap_drivers.png",
        ),
        (
            "Knowledge Graph",
            [
                f"Graph size: {metrics.kg_nodes:,} nodes and {metrics.kg_edges:,} edges.",
                f"Top critical node: {metrics.top_critical_node}.",
                "Relationships connect production lines, stations, features, transitions, and failure evidence.",
                "Failure propagation routes are monitoring hypotheses.",
            ],
            FIGURES_DIR / "phase9_critical_nodes.png",
        ),
        (
            "Manufacturing Copilot",
            [
                f"SQLite-backed copilot database contains {metrics.copilot_rows:,} prediction rows.",
                "Streamlit app supports risk review, root-cause Q&A, bottlenecks, graph insights, and summaries.",
                "Natural-language answers use deterministic reviewed intents and parameterized SQL.",
            ],
            None,
        ),
        (
            "Executive Dashboard",
            [
                "Dashboard includes KPI scorecard, relative failure trends, station heatmaps, bottlenecks, SHAP drivers, and business impact.",
                "Business impact is assumption-driven and explicitly labeled as a scenario.",
                "Dashboard is suitable for demos and future client adaptation after data-source replacement.",
            ],
            None,
        ),
        (
            "Deployment Boundary",
            [
                "Current system is a validated use-case prototype on Kaggle Bosch data.",
                "Future production use requires live plant data contracts, forward-time validation, access control, and monitoring.",
                "Model alerts should prioritize engineering review, not automate quality disposition without governance.",
            ],
            None,
        ),
    ]

    for title, bullets, image in slides:
        slide = prs.slides.add_slide(content_layout)
        add_slide_title(slide, title)
        add_bullets(slide, bullets, width=5.7 if image else 11.8)
        if image:
            add_picture(slide, image)

    slide = prs.slides.add_slide(content_layout)
    add_slide_title(slide, "Final Deliverable Map")
    add_bullets(
        slide,
        [
            "68 Model: models/phase6_best_model.joblib",
            "69 Segmentation: src/data/phase5_product_family_discovery.py",
            "70 Process mining: src/data/phase8_process_mining_bottleneck_analysis.py",
            "71 Root cause: src/data/phase7_root_cause_analysis.py",
            "72 Knowledge graph: reports/phase9_manufacturing_knowledge_graph.graphml",
            "73 Copilot: app/phase11_manufacturing_copilot.py",
            "74 Dashboard: app/phase12_executive_dashboard.py",
            "75 Documentation and presentation: docs/final_deliverables/",
        ],
        top=1.3,
        height=5.6,
    )

    path = DOCS_DIR / "Bosch_Production_Line_Performance_Presentation.pptx"
    prs.save(path)
    return path


def write_index(metrics: ProjectMetrics, markdown_path: Path, docx_path: Path, pptx_path: Path) -> Path:
    rows = deliverable_map()
    lines = [
        "# Final Deliverables Index",
        "",
        f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
        "",
        "## Headline Results",
        "",
        f"- Production-safe model: {metrics.best_model}",
        f"- Validation MCC: {metrics.mcc:.3f}",
        f"- Validation precision: {metrics.precision:.3f}",
        f"- Validation recall: {metrics.recall:.3f}",
        f"- Highest-risk station: {metrics.top_station} ({metrics.top_station_failure_rate_pct:.3f}% failure rate)",
        f"- Top bottleneck: {metrics.top_bottleneck} (score {metrics.top_bottleneck_score:.2f})",
        f"- Knowledge graph: {metrics.kg_nodes:,} nodes, {metrics.kg_edges:,} edges",
        "",
        "## Requested Deliverables",
        "",
        "| ID | Deliverable | Status | Primary artifacts |",
        "|---:|---|---|---|",
    ]
    for row in rows:
        artifacts = "<br>".join(f"`{artifact}`" for artifact in row["primary_artifacts"])
        lines.append(f"| {row['id']} | {row['deliverable']} | {row['status']} | {artifacts} |")
    lines.extend(
        [
            "",
            "## Documentation Artifacts",
            "",
            f"- Research Markdown: `{markdown_path.relative_to(PROJECT_ROOT)}`",
            f"- Research Word document: `{docx_path.relative_to(PROJECT_ROOT)}`",
            f"- Presentation deck: `{pptx_path.relative_to(PROJECT_ROOT)}`",
            "",
            "## Run Commands",
            "",
            "```powershell",
            ".\\.venv\\Scripts\\python.exe src\\data\\phase11_manufacturing_copilot.py",
            ".\\.venv\\Scripts\\python.exe src\\data\\phase12_executive_dashboard.py",
            ".\\.venv\\Scripts\\streamlit.exe run app\\phase11_manufacturing_copilot.py",
            ".\\.venv\\Scripts\\streamlit.exe run app\\phase12_executive_dashboard.py",
            "```",
        ]
    )
    path = DOCS_DIR / "final_deliverables_index.md"
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def write_manifest(markdown_path: Path, docx_path: Path, pptx_path: Path, index_path: Path) -> Path:
    manifest = {
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "project_root": str(PROJECT_ROOT),
        "deliverables": deliverable_map(),
        "documentation": {
            "research_markdown": str(markdown_path.relative_to(PROJECT_ROOT)),
            "research_docx": str(docx_path.relative_to(PROJECT_ROOT)),
            "presentation_pptx": str(pptx_path.relative_to(PROJECT_ROOT)),
            "index": str(index_path.relative_to(PROJECT_ROOT)),
        },
    }
    path = DOCS_DIR / "final_deliverables_manifest.json"
    path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    return path


def main() -> None:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    metrics = collect_metrics()
    markdown_path = write_markdown(metrics)
    docx_path = write_docx(metrics)
    pptx_path = write_presentation(metrics)
    index_path = write_index(metrics, markdown_path, docx_path, pptx_path)
    manifest_path = write_manifest(markdown_path, docx_path, pptx_path, index_path)
    print(
        json.dumps(
            {
                "markdown": str(markdown_path),
                "docx": str(docx_path),
                "pptx": str(pptx_path),
                "index": str(index_path),
                "manifest": str(manifest_path),
            },
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
